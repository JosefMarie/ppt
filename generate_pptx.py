import sys
import subprocess
import os

# 1. Self-install python-pptx if it's missing
try:
    import pptx
except ImportError:
    print("python-pptx library not found. Installing now...")
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        import pptx
    except Exception as e:
        print(f"Error installing python-pptx: {e}")
        print("Please install python-pptx manually: pip install python-pptx")
        sys.exit(1)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

def add_hyperlink(paragraph, text, link):
    """Add a hyperlink to a paragraph"""
    part = paragraph.part
    p = paragraph._element
    pPr = p.get_or_add_pPr()
    
    # Create hyperlink run
    run = paragraph.add_run()
    run.text = text
    
    # Add hyperlink relationship
    rId = part.relate_to(link, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True)
    
    # Create the hyperlink element
    hyperlink = OxmlElement('a:hyperlink')
    hyperlink.set('r:id', rId)
    run._r.addprevious(hyperlink)
    run._r.getparent().remove(run._r)
    hyperlink.append(run._r)
    
def create_pitch_deck():
    prs = Presentation()
    
    # Set to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom Brand Colors
    bg_color = RGBColor(5, 5, 5)          # Deep Black #050505
    cyan_color = RGBColor(0, 242, 255)    # Electric Cyan #00f2ff
    purple_color = RGBColor(168, 85, 247) # Cyber Purple #a855f7
    white_color = RGBColor(249, 250, 251) # Soft White #f9fafb
    grey_color = RGBColor(156, 163, 175)  # Muted Grey #9ca3af
    
    # Blank layout is layout index 6
    blank_layout = prs.slide_layouts[6]
    
    # Slide data blueprint
    slides_data = [
        # Slide 1
        {
            "category": "ENGINEERING THE NEXT 2026",
            "title": "DIGITAL\nFOUNDATION",
            "subtitle": "Engineering the Digital Future of East Africa",
            "footer": "Presenter: Ashimirwe Joseph Marie  |  Kigali, Rwanda",
            "bullets": [],
            "notes": "I’m Ashimirwe Joseph Marie. We are Foundation[EA], a premium digital engineering agency based in Kigali, dedicated to setting a new aesthetic and functional standard for software in our region."
        },
        # Slide 2
        {
            "category": "THE \"WHY\"",
            "title": "Bridging the Digital Divide",
            "subtitle": "",
            "bullets": [
                "The East African digital landscape is growing rapidly.",
                "Businesses need more than just code; they need high-end, future-proof digital infrastructure.",
                "We bridge the gap between classroom theory and enterprise-grade reality."
            ],
            "notes": "We see a shift toward digitalization, but many businesses are stuck with outdated, analog processes. We exist to provide the high-end, '2026-Aesthetic' digital solutions they need to compete globally."
        },
        # Slide 3
        {
            "category": "THE OPPORTUNITY",
            "title": "The \"Analog\" Friction",
            "subtitle": "",
            "bullets": [
                "Digital solutions often lack intuitive design and performance.",
                "Enterprises struggle with inefficient, paper-heavy workflows.",
                "A local talent gap exists between academic learning and industry standards."
            ],
            "notes": "Digital solutions often lack intuitive design and performance. Enterprises struggle with inefficient, paper-heavy workflows. A local talent gap exists between academic learning and industry standards."
        },
        # Slide 4
        {
            "category": "OUR APPROACH",
            "title": "Premium Engineering, Local Impact",
            "subtitle": "",
            "bullets": [
                "Expert Stack: Next.js 15, TypeScript, Python, Node.js, Firebase, PostgreSQL.",
                "Signature Aesthetic: Kinetic typography, glassmorphism, and real-time synchronization.",
                "Focus: Building software that is as beautiful as it is functional."
            ],
            "notes": "We build premium solutions using Next.js 15, TypeScript, Python, Node.js, and modern databases. We pair this with a rich visual style featuring kinetic typography, glassmorphism, and real-time syncing."
        },
        # Slide 5
        {
            "category": "BUSINESS ENGINE",
            "title": "Engineering & Incubation",
            "subtitle": "",
            "bullets": [
                "Professional Agency: High-end digital engineering for enterprise clients.",
                "Talent Incubation: Bridging the gap for local developers through mentorship and real-world projects.",
                "Founder Synergies: Maintaining a pipeline of top-tier talent while delivering world-class service."
            ],
            "notes": "My dual role as a founder and a software development teacher allows us to maintain a pipeline of top-tier talent while delivering world-class service."
        },
        # Slide 6
        {
            "category": "OUR WORK",
            "title": "Our Flagship Work",
            "subtitle": "",
            "bullets": [
                "Foundation[EA]: Premium digital engineering with a local-first, global-standard approach.",
                "USAFI Barista Training Center: Professional coffee training in brewing, espresso, latte art, and specialty drinks.",
                "Divine's Destinations: Refined travel storytelling with cinematic UX and precise digital interactions.",
                "Happy Fish Farm Rwanda Academy: Sustainable aquaculture education for future-ready farm leaders."
            ],
            "notes": "Our flagship portfolio now spans digital engineering, hospitality training, travel, aquaculture education, and precision agriculture. Each product is designed to solve real local problems with modern technology."
        },
        # Slide 7
        {
            "category": "OUR WORK",
            "title": "Our Flagship Work (Continued)",
            "subtitle": "",
            "bullets": [
                "5A FARMS: Precision agriculture and modular ecosystems for resilient food systems.",
                "Portfolio Focus: Education, hospitality, travel, and agriculture—each built to create measurable local impact."
            ],
            "notes": "We continue to expand our portfolio through practical, high-impact experiences that connect education, industry, and sustainable growth across East Africa."
        },
        # Slide 8
        {
            "category": "THE HORIZON",
            "title": "The Digital Horizon",
            "subtitle": "",
            "bullets": [
                "Kigali is emerging as a regional tech hub.",
                "Rising demand for digital transformation in education, logistics, and tourism sectors.",
                "Foundation[EA] is positioned at the intersection of this regional ICT vision."
            ],
            "notes": "Kigali is emerging as a regional tech hub. Rising demand for digital transformation in education, logistics, and tourism sectors. Foundation[EA] is positioned at the intersection of this regional ICT vision."
        },
        # Slide 9
        {
            "category": "GROWTH STRATEGY",
            "title": "Scaling Our Velocity (10M RWF)",
            "subtitle": "",
            "bullets": [
                "3.5M RWF: High-performance hardware (Laptops/Workstations).",
                "2.5M RWF: Establishing a centralized Studio HQ in Kigali.",
                "1.5M RWF: Cloud infrastructure, enterprise API licenses, and dev tooling.",
                "1.0M RWF: Branding, marketing, and professional portfolio development.",
                "1.5M RWF: Innovation/R&D buffer for market-shifting prototypes."
            ],
            "notes": "We are seeking a 10M RWF growth runway: 3.5M for hardware, 2.5M for a Studio HQ in Kigali, 1.5M for cloud and API tooling, 1.0M for portfolio branding, and 1.5M for an R&D prototype buffer."
        },
        # Slide 9
        {
            "category": "LEGACY",
            "title": "Building a Legacy",
            "subtitle": "",
            "bullets": [
                "Professionalizing the regional software engineering ecosystem.",
                "Empowering students to solve local problems with global-standard technology.",
                "Creating a permanent, innovation-focused infrastructure in Kigali."
            ],
            "notes": "We are building a legacy: professionalizing the regional software engineering ecosystem, empowering students to solve local problems, and creating permanent innovation infrastructure."
        },
        # Slide 10
        {
            "category": "JOIN US",
            "title": "Join Us",
            "subtitle": "Let's build the digital backbone of East Africa together.",
            "bullets": [
                "Website: foundationea.com",
                "📧 Email: agency@mail.foundationea.com",
                "💬 WhatsApp: +250783309973",
                "Founder: Ashimirwe Joseph Marie"
            ],
            "notes": "I invite you to join us in building this infrastructure. Let's make Foundation[EA] the premier name in African digital engineering. Contact us via WhatsApp at +250783309973 or email agency@mail.foundationea.com."
        }
    ]
    
    for i, data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        
        # Apply dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # Add a subtle grid effect or accent border
        # We draw a thin accent colored shape on the left border of the slide
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), Inches(0.15), Inches(7.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = cyan_color if i % 2 == 0 else purple_color
        accent_bar.line.color.rgb = cyan_color if i % 2 == 0 else purple_color
        
        # Speaker notes configuration
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = data["notes"]
        
        # Category label textbox (e.g. THE "WHY")
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = data["category"]
        cat_p.font.name = "Arial"
        cat_p.font.size = Pt(11)
        cat_p.font.bold = True
        cat_p.font.color.rgb = cyan_color
        
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(1.3))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_p = title_tf.paragraphs[0]
        title_p.text = data["title"]
        title_p.font.name = "Arial"
        # Large headers for titles
        title_p.font.size = Pt(44) if i == 0 else Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = white_color if i != 0 else purple_color
        
        # Slide Title 1 custom alignment
        if i == 0:
            title_p.alignment = PP_ALIGN.LEFT
            title_box.top = Inches(1.8)
            title_box.height = Inches(2.2)
        
        # Subtitle (if present)
        if data["subtitle"]:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8) if i == 0 else Inches(2.2), Inches(11.7), Inches(0.6))
            sub_tf = sub_box.text_frame
            sub_tf.word_wrap = True
            sub_p = sub_tf.paragraphs[0]
            sub_p.text = data["subtitle"]
            sub_p.font.name = "Arial"
            sub_p.font.size = Pt(18)
            sub_p.font.color.rgb = grey_color
            
        # Bullets / Core Content List
        if data["bullets"]:
            bullet_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.7), Inches(4.0))
            bullet_tf = bullet_box.text_frame
            bullet_tf.word_wrap = True
            
            for index, bullet in enumerate(data["bullets"]):
                if index == 0:
                    p = bullet_tf.paragraphs[0]
                else:
                    p = bullet_tf.add_paragraph()
                
                p.text = "•  " + bullet
                p.space_after = Pt(20)
                p.font.name = "Arial"
                p.font.size = Pt(20)
                p.font.color.rgb = white_color
                
        # Footer branding for slide 1
        if i == 0 and data["footer"]:
            foot_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.5))
            foot_tf = foot_box.text_frame
            foot_p = foot_tf.paragraphs[0]
            foot_p.text = data["footer"]
            foot_p.font.name = "Arial"
            foot_p.font.size = Pt(11)
            foot_p.font.color.rgb = grey_color

    # Save presentation
    output_filename = "FoundationEA_Pitch_Deck.pptx"
    prs.save(output_filename)
    print(f"PowerPoint Presentation successfully created: {os.path.abspath(output_filename)}")

if __name__ == "__main__":
    create_pitch_deck()
